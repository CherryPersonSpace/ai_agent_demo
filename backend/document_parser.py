"""
文档内容提取模块
支持格式: .txt .docx .xls .xlsx .pptx .pdf .png .jpg .jpeg .bmp .tiff .tif .webp
图片文件通过 RapidOCR 进行文字识别（OCR）。
"""
import os
import io
from pathlib import Path

SUPPORTED_EXTENSIONS = {
    ".txt", ".docx", ".xls", ".xlsx", ".pptx", ".pdf",
    ".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp",
}

# 图片文件后缀集合（用于 OCR 识别）
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp"}


# ────────────────────────────────────────────
# 1. TXT
# ────────────────────────────────────────────
def extract_text_from_txt(file_path: str) -> str:
    """直接读取纯文本文件，自动检测编码。"""
    # 尝试 utf-8，失败回退 gbk（常见于中文 Windows）
    for enc in ("utf-8", "gbk", "latin-1"):
        try:
            with open(file_path, "r", encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    raise ValueError(f"无法解码文本文件: {file_path}")


# ────────────────────────────────────────────
# 2. DOCX（复用已有的解析逻辑）
# ────────────────────────────────────────────
def extract_text_from_docx(file_path: str) -> str:
    """从 .docx 文件中提取纯文本（含表格）。"""
    from docx import Document
    from docx.oxml.ns import qn
    from docx.text.paragraph import Paragraph
    from docx.table import Table

    doc = Document(file_path)
    texts = []

    for child in doc.element.body:
        if child.tag == qn("w:p"):
            para = Paragraph(child, doc)
            text = para.text.strip()
            if text:
                texts.append(text)
        elif child.tag == qn("w:tbl"):
            table = Table(child, doc)
            table_lines = []
            for row_idx, row in enumerate(table.rows):
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                row_text = " | ".join(cells)
                if row_text.strip():
                    table_lines.append(row_text)
                if row_idx == 0 and table_lines:
                    table_lines.append(" | ".join(["---"] * len(cells)))
            if table_lines:
                texts.append("[表格]\n" + "\n".join(table_lines))

    return "\n\n".join(texts)


# ────────────────────────────────────────────
# 3. XLS / XLSX
# ────────────────────────────────────────────
def extract_text_from_excel(file_path: str) -> str:
    """从 Excel 文件中提取所有工作表的文本内容。"""
    ext = Path(file_path).suffix.lower()

    if ext == ".xls":
        # 旧版二进制格式用 xlrd
        import xlrd
        wb = xlrd.open_workbook(file_path)
        sheets_text = []
        for sheet in wb.sheets():
            rows = []
            for row_idx in range(sheet.nrows):
                cells = [str(sheet.cell_value(row_idx, col)).strip()
                         for col in range(sheet.ncols)]
                row_text = " | ".join(cells)
                if row_text.strip(" |"):
                    rows.append(row_text)
            if rows:
                sheets_text.append(f"[工作表: {sheet.name}]\n" + "\n".join(rows))
        return "\n\n".join(sheets_text)
    else:
        # .xlsx 用 openpyxl，失败则回退 xlrd（可能是旧格式 .xls 被重命名为 .xlsx）
        import zipfile
        try:
            import openpyxl
        except ImportError:
            raise ValueError(
                "缺少 openpyxl 库，无法解析 .xlsx 文件。"
                "请在虚拟环境中执行: pip install openpyxl"
            )

        # 检查文件大小
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            raise ValueError(f"文件 '{Path(file_path).name}' 为空文件")

        # 检查文件头（xlsx 是 zip 格式，头两个字节应为 PK）
        with open(file_path, "rb") as f:
            header = f.read(4)

        if header[:2] != b"PK":
            # 不是 zip 格式，可能是旧版 .xls 被重命名为 .xlsx
            try:
                import xlrd
                wb_xls = xlrd.open_workbook(file_path)
                sheets_text = []
                for sheet in wb_xls.sheets():
                    rows = []
                    for row_idx in range(sheet.nrows):
                        cells = [str(sheet.cell_value(row_idx, col)).strip()
                                 for col in range(sheet.ncols)]
                        row_text = " | ".join(cells)
                        if row_text.strip(" |"):
                            rows.append(row_text)
                    if rows:
                        sheets_text.append(f"[工作表: {sheet.name}]\n" + "\n".join(rows))
                return "\n\n".join(sheets_text)
            except Exception:
                raise ValueError(
                    f"文件 '{Path(file_path).name}' 不是有效的 Excel 文件"
                    f"（文件头: {header[:4].hex()}，大小: {file_size} 字节）"
                )

        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
        except Exception as e:
            err_type = type(e).__name__
            err_msg = str(e)
            # 列出 zip 内容以便诊断
            zip_info = ""
            try:
                with zipfile.ZipFile(file_path, "r") as zf:
                    zip_info = f"\nzip 内容: {zf.namelist()[:10]}"
            except Exception:
                zip_info = "\n（无法读取 zip 内容）"
            raise ValueError(
                f"解析 .xlsx 失败 [{err_type}]: {err_msg}"
                f"{zip_info}\n文件: {Path(file_path).name}（{file_size} 字节，头: {header[:4].hex()}）"
            )

        sheets_text = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                row_text = " | ".join(cells)
                if row_text.strip(" |"):
                    rows.append(row_text)
            if rows:
                sheets_text.append(f"[工作表: {sheet.title}]\n" + "\n".join(rows))
        return "\n\n".join(sheets_text)


# ────────────────────────────────────────────
# 4. PPTX
# ────────────────────────────────────────────
def extract_text_from_pptx(file_path: str) -> str:
    """从 PowerPoint (.pptx) 文件中提取所有幻灯片的文本。"""
    from pptx import Presentation

    prs = Presentation(file_path)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            # 表格
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    row_text = " | ".join(cells)
                    if row_text.strip(" |"):
                        texts.append(row_text)
        if texts:
            slides_text.append(f"[幻灯片 {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides_text)


# ────────────────────────────────────────────
# 5. 图片 OCR
# ────────────────────────────────────────────
def extract_text_from_image(file_path: str) -> str:
    """使用 RapidOCR 从图片中提取文字内容。"""
    try:
        from rapidocr_onnxruntime import RapidOCR
    except ImportError:
        raise ValueError(
            "缺少 rapidocr_onnxruntime 库，无法识别图片中的文字。"
            "请在虚拟环境中执行: pip install rapidocr_onnxruntime"
        )

    ocr = RapidOCR()
    result, _ = ocr(file_path)

    if not result:
        raise ValueError(
            f"图片 '{Path(file_path).name}' 中未识别到任何文字内容。"
            "请确认图片清晰且包含可读文字。"
        )

    # result: List[List[Any]]，每个元素为 [bbox, text, confidence]
    lines = [item[1] for item in result]
    return "\n".join(lines)


# ────────────────────────────────────────────
# 6. PDF
# ────────────────────────────────────────────
def extract_text_from_pdf(file_path: str) -> str:
    """从 PDF 文件中提取各页的文本内容。"""
    import fitz  # PyMuPDF

    doc = fitz.open(file_path)
    pages_text = []
    for i, page in enumerate(doc, 1):
        text = page.get_text().strip()
        if text:
            pages_text.append(f"[第 {i} 页]\n{text}")
    doc.close()
    return "\n\n".join(pages_text)


# ────────────────────────────────────────────
# 统一入口
# ────────────────────────────────────────────
_EXTRACTORS = {
    ".txt": extract_text_from_txt,
    ".docx": extract_text_from_docx,
    ".xls": extract_text_from_excel,
    ".xlsx": extract_text_from_excel,
    ".pptx": extract_text_from_pptx,
    ".pdf": extract_text_from_pdf,
    ".png": extract_text_from_image,
    ".jpg": extract_text_from_image,
    ".jpeg": extract_text_from_image,
    ".bmp": extract_text_from_image,
    ".tiff": extract_text_from_image,
    ".tif": extract_text_from_image,
    ".webp": extract_text_from_image,
}


def extract_document(file_path: str) -> str:
    """
    根据文件后缀自动选择合适的解析器，提取纯文本内容。

    Args:
        file_path: 文件的本地路径

    Returns:
        提取到的纯文本内容

    Raises:
        ValueError: 不支持的文件格式
        Exception: 解析过程中可能出现的其他错误
    """
    ext = Path(file_path).suffix.lower()
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(f"不支持的文件格式 '{ext}'，目前支持: {supported}")
    return extractor(file_path)
